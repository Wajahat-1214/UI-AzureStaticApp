import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Create directories to store outputs
OUTPUT_DIR = "pptx_extraction"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Output file path for saving extracted data
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "extracted_pptx_data.txt")

# Load the presentation
PPTX_FILE = "Postman_presentation.pptx"
prs = Presentation(PPTX_FILE)

def write_to_file(content):
    """Helper function to write content to the output file."""
    with open(OUTPUT_FILE, "a", encoding="utf-8") as f:
        f.write(content + "\n")

def extract_text_and_fonts():
    write_to_file("\n--- Extracting Text and Font Properties ---\n")
    for i, slide in enumerate(prs.slides):
        write_to_file(f"Slide {i + 1}:")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text
                    write_to_file(f"  Paragraph: {text}")
                    for run in paragraph.runs:
                        font = run.font
                        font_size = font.size.pt if font.size else "Default"
                        
                        # Safely handle font color
                        font_color = "Default"
                        if font.color and font.color.type:  # Check if a color type exists
                            if font.color.type == 1 and font.color.rgb:  # Solid color type
                                font_color = font.color.rgb
                        
                        bold = font.bold if font.bold is not None else "Default"
                        write_to_file(f"    Font Size: {font_size}")
                        write_to_file(f"    Font Color: {font_color}")
                        write_to_file(f"    Bold: {bold}")



def extract_tables():
    write_to_file("\n--- Extracting Tables ---\n")
    for i, slide in enumerate(prs.slides):
        write_to_file(f"Slide {i + 1}:")
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    row_data = [cell.text for cell in row.cells]
                    write_to_file(f"  Row {row_idx + 1}: {row_data}")

def extract_images():
    write_to_file("\n--- Extracting Images ---\n")
    img_count = 0
    img_dir = os.path.join(OUTPUT_DIR, "images")
    os.makedirs(img_dir, exist_ok=True)
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                img_bytes = image.blob
                img_path = os.path.join(img_dir, f"image_{img_count}.png")
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                write_to_file(f"Extracted image {img_count} from slide {i + 1} to {img_path}")
                img_count += 1

def extract_layouts():
    write_to_file("\n--- Extracting Slide Layouts ---\n")
    for i, layout in enumerate(prs.slide_layouts):
        write_to_file(f"Layout {i + 1}: {layout.name}")

def main():
    # Clear the output file before writing
    open(OUTPUT_FILE, "w").close()
    
    write_to_file("Starting PowerPoint Extraction...\n")
    extract_text_and_fonts()
    extract_tables()
    extract_images()
    extract_layouts()
    write_to_file("\nExtraction Complete!")

if __name__ == "__main__":
    main()
